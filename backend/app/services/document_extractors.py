from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pandas as pd
from docx import Document
from pypdf import PdfReader
from pptx import Presentation


def extract_text(filename: str, content: bytes) -> str:
    suffix = Path(filename).suffix.lower()

    if suffix == '.pdf':
        reader = PdfReader(BytesIO(content))
        return '\n'.join(page.extract_text() or '' for page in reader.pages)

    if suffix == '.docx':
        document = Document(BytesIO(content))
        return '\n'.join(paragraph.text for paragraph in document.paragraphs)

    if suffix == '.txt':
        return content.decode('utf-8', errors='ignore')

    if suffix == '.csv':
        return pd.read_csv(BytesIO(content)).to_csv(index=False)

    if suffix == '.pptx':
        presentation = Presentation(BytesIO(content))
        blocks: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    blocks.append(shape.text)
        return '\n'.join(blocks)

    return content.decode('utf-8', errors='ignore')

